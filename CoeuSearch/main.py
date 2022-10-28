#---------------------------------------------------------------------------------------------------------------------------
# IMPORTS
#---------------------------------------------------------------------------------------------------------------------------
import os
import csv
import docx
import pptx
import time
import nltk
import pickle
import PyPDF2
import numpy as np
import pandas as pd
import configs
from tqdm import tqdm
from keybert import KeyBERT
from nltk.stem import WordNetLemmatizer 
from nltk.tokenize import sent_tokenize, word_tokenize
from sentence_transformers import SentenceTransformer
from CoeuSearch.utils import parser_pdf, cleanText, Color
from sklearn.metrics.pairwise import cosine_similarity, euclidean_distances

#---------------------------------------------------------------------------------------------------------------------------
# CONFIGS
#---------------------------------------------------------------------------------------------------------------------------
# nltk.download('all')
color = Color()
lemmatizer = WordNetLemmatizer()
use_cache = configs.use_cache
model_name = configs.model_name

kw_model = KeyBERT()
sbert_model = SentenceTransformer(model_name)

cache_base_path = configs.cache_base_path
if not os.path.exists(cache_base_path):
    os.makedirs(cache_base_path)

#---------------------------------------------------------------------------------------------------------------------------
# FILE OBJECT
#---------------------------------------------------------------------------------------------------------------------------
class File:
  def __init__(self, path, ftype):
    self.path = path
    self.ftype = ftype #.txt .docx .doc .pdf
    self.name = None
    self.title = None
    self.content = None 
    self.title_embed = None
    self.content_embed = None
    self.content_cache_path = None
    self.title_embed_cache_path = None
    self.content_embed_cache_path = None
    
    def setFileConfig(self):
        self.name = ".".join(self.path.split('\\')[-1].split('.')[:-1])
        temp_path = self.path.replace("\\", "-").replace(":", "-").replace(".", "-")
        self.content_cache_path = cache_base_path + temp_path +"__content.txt"
        self.title_embed_cache_path = cache_base_path + temp_path +"__title.npy"
        self.content_embed_cache_path = cache_base_path + temp_path +"__content.npy"
    
    def setTitle(self):
        temp_name = self.name.replace('-', ' ').replace('_', ' ')
        temp_name_words = word_tokenize(temp_name)
        self.title = " ".join([lemmatizer.lemmatize(w) for w in temp_name_words])
    
    def cacheCheck(self):
        #CHECKING TITLE EMBEDS
        if os.path.exists(self.title_embed_cache_path):
            self.title_embed = np.load(self.title_embed_cache_path)
        else:
            return 0
        #CHECKING CONTENT EMBEDS
        if os.path.exists(self.content_embed_cache_path):
            self.content_embed = np.load(self.content_embed_cache_path)
        else:
            return 0
        #CHECKING CONTENT STRING
        if os.path.exists(self.content_cache_path):
            with open(self.content_cache_path, "r", encoding="utf-8") as f:
                self.content = f.read()
        else:
            return 0        
        return 1
        

    def readContent(self):
        lines = []
        # [PDF FILE] READING CONTENT
        if ftype == "pdf":
            try:
                pdf_text = parser_pdf(self.path)
                lines = pdf_text.split("\n")
            except Exception as E:
                print( color.RED+ f"Error in reading pdf file: {self.path} - {E}" + color.END)
        # [DOC FILE] READING CONTENT
        if ftype == 'doc':
            try:
                with open(self.path, "rb") as f:
                    doc = docx.Document(self.path)
                    lines = [str(para.text) for para in doc.paragraphs]
            except Exception as E:
                print( color.RED+ f"Error in reading docx file: {self.path} - {E}" + color.END)
        # [PPT FILE] READING CONTENT
        if ftype == "ppt":
            try:
                text = ""
                prs = pptx.Presentation(self.path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = text + str(shape.text) + " \n"
                lines = text.split("\n")
            except Exception as E:
                print( color.RED+ f"Error in reading pptx file: {self.path} - {E}" + color.END)
        # [csv FILE] READING CONTENT
        if ftype == "csv":
            text =[]
            try:
                with open( self.path, 'rb') as csvfile:
                    reader_ = csv.reader(csvfile, delimiter=' ', quotechar='|')
                    for row in reader_:
                        text = text + ' '.join(row)+ " \n"
                lines = text.split("\n")
            except Exception as E:
                print( color.RED+ f"Error in reading pptx file: {self.path} - {E}" + color.END)
        # [TXT FILE] READING CONTENT
        if ftype == 'txt':
            try:
                with open(self.path) as f:
                    lines = f.readlines()
                f.close()
            except Exception as E:
                print( color.RED+ f"Error in reading txt file: {self.path} - {E}" + color.END)
        
        self.content = cleanText(lines)
        if self.content:
            with open(self.content_cache_path, "w", encoding="utf-8") as g:
                g.write(self.content)
            g.close()

    def processEmbeds(self):
        # TITLE EMBEDS
        if self.title:
            self.title_embed = sbert_model.encode([self.title])[0].reshape(1,-1)
            np.save(self.title_embed_cache_path, self.title_embed)
        # CONTENT EMBEDS
        if self.content:
            self.content_embed = sbert_model.encode([self.content])[0].reshape(1,-1)
            np.save(self.content_embed_cache_path, self.content_embed)
            

    setFileConfig(self)
    setTitle(self)

    if use_cache and cacheCheck(self):
        print( color.CYAN + f"\t[LOADING DATA FROM CACHE]: {self.name}" + color.END)
    else:
        print( color.CYAN + f"\t[PREPARING NEW DATA]: {self.name}" + color.END)
        readContent(self)
        processEmbeds(self)

#---------------------------------------------------------------------------------------------------------------------------
# GET FILES IN DIRECTORY | PROCESS QUERY | CALL SEARCH
#---------------------------------------------------------------------------------------------------------------------------
def getFiles(path, query, probs):
    if (path) and (query):
        #RECURSIVELY GET PATHS
        print( color.GREEN + "[LOADING FILES]" + color.END)
        txt_file_paths, pdf_file_paths, doc_file_paths, ppt_file_paths, csv_file_paths, image_file_paths, other_file_paths = [],[],[],[],[],[],[]
        try:
            for r, d, f in os.walk(path):
                for file in f:
                    if '.txt' in file:
                        txt_file_paths.append(os.path.join(r, file))
                    elif '.pdf' in file:
                        pdf_file_paths.append(os.path.join(r, file))
                    elif '.docx' in file or '.doc' in file:
                        doc_file_paths.append(os.path.join(r, file))
                    elif '.pptx' in file or '.ppt' in file:
                        ppt_file_paths.append(os.path.join(r, file))
                    elif '.csv' in file or '.xlsx' in file or '.xls' in file:
                        csv_file_paths.append(os.path.join(r, file))
                    elif '.jpg' in file or '.jpeg' in file or '.png' in file:
                        image_file_paths.append(os.path.join(r, file))
                    else:
                        other_file_paths.append(os.path.join(r, file))
        except Exception as E:
            return {'error': "Invalid Path: " + str(E)}
        
        #CREATE FILE OBJECTS
        print( color.GREEN + "[CREATING FILE OBJECTS]" + color.END)
        txt_files = [File(path, "txt") for path in txt_file_paths]
        pdf_files = [File(path, "pdf") for path in pdf_file_paths]
        doc_files = [File(path, "doc") for path in doc_file_paths]
        ppt_files = [File(path, "ppt") for path in ppt_file_paths]
        csv_files = [File(path, "csv") for path in csv_file_paths]
        image_files = [File(path, "image") for path in image_file_paths]
        other_files = [File(path, "other") for path in other_file_paths]
        all_files = txt_files + pdf_files + doc_files + ppt_files + csv_files + image_files + other_files
        
        #PREPROCESS QUERY
        print( color.GREEN + "[PROCESSING QUERY]" + color.END)
        temp_query_words = word_tokenize(query.replace('-', ' ').replace('_', ' '))
        proc_query = " ".join([lemmatizer.lemmatize(w) for w in temp_query_words])

        #SEARCH
        print( color.GREEN + "[SEARCHING]" + color.END)
        relevant_files = search(all_files, proc_query, probs)

        return {'files':[ (x.path, x.ftype) for x in relevant_files]},  time.time()
    else:
        return None


#---------------------------------------------------------------------------------------------------------------------------
# NEURAL SEARCH
#---------------------------------------------------------------------------------------------------------------------------
def search(files, proc_query, probs):
    query_tokens = len(proc_query.split(' '))
    query_embed = sbert_model.encode([proc_query])[0].reshape(1,-1)
    
    # MODEL 1: TITLE MATCHING
    print( color.CYAN + "\t[MODEL 1: MATCHING TTILES WITH QUERY]" + color.END)
    title_scores = []
    for f in files:
        key_embed = f.title_embed.reshape(1,-1)
        title_scores.append(cosine_similarity(key_embed, query_embed)[0][0])
        
    # MODEL 2: CONTENT MATCHING
    print( color.CYAN + "\t[MODEL 2: MATCHING CONTENT WITH QUERY]" + color.END)
    content_scores = []
    for f in files:
        if f.content_embed is not None:
            key_embed = f.content_embed.reshape(1,-1)
            content_scores.append(cosine_similarity(key_embed, query_embed)[0][0])
        else:
            content_scores.append(0.0)            

    # MODEL 3: TOPIC MATCHING
    print( color.CYAN + "\t[MODEL 3: MATCHING CONTENT TOPICS WITH QUERY]" + color.END)
    topics, topic_scores = [], []
    for i, f in tqdm(enumerate(files)):
        if f.content is not None:
            content_keywords_embed_cache_path = f.content_embed_cache_path[:-4] + f"qt_{query_tokens}_keywords.txt"
            if os.path.exists(content_keywords_embed_cache_path):
                with open(content_keywords_embed_cache_path, "rb") as fp:
                    top_keywords = pickle.load(fp)       
            else:
                keyword_and_score = kw_model.extract_keywords(f.content, keyphrase_ngram_range=(1, query_tokens), stop_words=None)
                top_keywords = [kws[0] for kws in keyword_and_score]                
                with open(content_keywords_embed_cache_path, "wb") as fp:
                    pickle.dump(top_keywords, fp) 
            top_keywords_embed = sbert_model.encode(top_keywords)
            keyword_scores = [cosine_similarity(k_embed.reshape(1,-1), query_embed)[0][0] for k_embed in top_keywords_embed] + [0.0]
            best_kw_id = np.argmax(keyword_scores)        
            if top_keywords:
                topics.append(top_keywords[best_kw_id])
            else:
                topics.append(None)
            topic_scores.append(keyword_scores[best_kw_id])
        else:
            topic_scores.append(0.0)
            topics.append(None)

    # ENSEMBLE MODEL: COMBINE RESULTS
    results = pd.DataFrame({'file': files, 
                            'name': [f.name for f in files],
                            'ftype': [f.ftype for f in files],
                            'content': [f.content for f in files],
                            'topic': topics, 
                            'title_score': title_scores,
                            'content_score': content_scores,
                            'topic_score': topic_scores,
                           }).sort_values(by=['title_score'], ascending=False)

    title_results = results[results['title_score'] > probs['title_prob_thres']]
    content_results = results[results['content_score'] > probs['content_prob_thres']]
    topic_results = results[results['topic_score'] > probs['topic_prob_thres']]
    relevant_files = list(set(list(title_results['file']) + list(content_results['file']) + list(topic_results['file'])))
    
    display_df = results[['name', 'ftype', 'topic', 'title_score', 'content_score', 'topic_score']]
    display_df.to_csv("results.csv")
    
    print( color.GREEN + "[RENDERING RESULTS]\n" + color.END, display_df)
    return relevant_files


#---------------------------------------------------------------------------------------------------------------------------
# END
#---------------------------------------------------------------------------------------------------------------------------